# logic/analyzer/ingest.py
from __future__ import annotations
import os
import io
import base64
import hashlib
from typing import List, Dict, Any, Tuple

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image

from config import MAX_IMGS_PER_PAGE, MIN_IMAGE_AREA, DEDUPLICATE_IMAGES


def _b64_data_url(img_bytes: bytes, prefer_format: str | None = None) -> Tuple[str, int, int, int]:
    """
    Convierte una imagen a data URL base64, reescala si el mayor lado >1600 px.
    Devuelve (data_url, width, height, area).
    """
    with Image.open(io.BytesIO(img_bytes)) as im:
        im = im.convert("RGB")
        w, h = im.size
        max_side = max(w, h)
        if max_side > 1600:
            scale = 1600.0 / max_side
            im = im.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        buf = io.BytesIO()
        out_fmt = "JPEG"
        if prefer_format:
            pf = prefer_format.lower()
            if pf == "png":
                out_fmt = "PNG"
        im.save(buf, format=out_fmt, optimize=True, quality=85)
        data = base64.b64encode(buf.getvalue()).decode("utf-8")
        media = "image/png" if out_fmt == "PNG" else "image/jpeg"
        w2, h2 = im.size
        return f"data:{media};base64,{data}", w2, h2, (w2 * h2)


def _hash_bytes(data: bytes) -> str:
    return hashlib.md5(data).hexdigest()


def ingest_files(paths: List[str]) -> List[Dict[str, Any]]:
    """
    Ingesta unificada para PDF, DOCX y PPTX. Devuelve items:
      - {"kind":"text","text":..., "page":N}
      - {"kind":"image","data_url":..., "page":N, "width":W, "height":H, "area":A, "caption": None}
    No aplica OCR aquí (se hace después para mantener esta fase libre de API externas).
    """
    items: List[Dict[str, Any]] = []
    seen_hashes: set[str] = set()

    for path in paths:
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pdf":
            items.extend(_ingest_pdf(path, seen_hashes))
        elif ext == ".docx":
            items.extend(_ingest_docx(path, seen_hashes))
        elif ext == ".pptx":
            items.extend(_ingest_pptx(path, seen_hashes))
        else:
            continue
    return items


def _ingest_pdf(path: str, seen_hashes: set[str]) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    doc = fitz.open(path)
    try:
        for page_idx, page in enumerate(doc, start=1):
            text = page.get_text("text") or ""
            if text.strip():
                out.append({"kind": "text", "text": text, "page": page_idx})

            page_imgs: List[Tuple[int, bytes, str]] = []  # (area, raw_bytes, ext)
            for img in page.get_images(full=True):
                try:
                    xref = img[0]
                    img_dict = doc.extract_image(xref)
                    img_bytes = img_dict["image"]
                    hsh = _hash_bytes(img_bytes)
                    if DEDUPLICATE_IMAGES and hsh in seen_hashes:
                        continue
                    with Image.open(io.BytesIO(img_bytes)) as im:
                        w, h = im.size
                        area = w * h
                    if area < MIN_IMAGE_AREA:
                        continue
                    page_imgs.append((area, img_bytes, (img_dict.get("ext") or "jpeg").lower()))
                    if DEDUPLICATE_IMAGES:
                        seen_hashes.add(hsh)
                except Exception:
                    continue

            if page_imgs:
                page_imgs.sort(key=lambda x: x[0], reverse=True)
                for area, raw, ext in page_imgs[:MAX_IMGS_PER_PAGE]:
                    data_url, w2, h2, area2 = _b64_data_url(raw, prefer_format=ext)
                    out.append({
                        "kind": "image",
                        "data_url": data_url,
                        "page": page_idx,
                        "width": w2,
                        "height": h2,
                        "area": area2,
                        "caption": None
                    })
    finally:
        doc.close()
    return out


def _ingest_docx(path: str, seen_hashes: set[str]) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    doc = DocxDocument(path)

    texts = []
    for p in doc.paragraphs:
        if p.text and p.text.strip():
            texts.append(p.text)
    full_text = "\n".join(texts).strip()
    if full_text:
        out.append({"kind": "text", "text": full_text, "page": 1})

    try:
        rels = doc.part._rels
        page_imgs: List[Tuple[int, bytes, str]] = []
        for rel in rels.values():
            if "image" in rel.target_ref or ("/media/" in rel.target_ref):
                part = rel._target
                blob: bytes = part.blob
                hsh = _hash_bytes(blob)
                if DEDUPLICATE_IMAGES and hsh in seen_hashes:
                    continue
                with Image.open(io.BytesIO(blob)) as im:
                    w, h = im.size
                    area = w * h
                if area < MIN_IMAGE_AREA:
                    continue
                page_imgs.append((area, blob, "jpeg"))
                if DEDUPLICATE_IMAGES:
                    seen_hashes.add(hsh)

        if page_imgs:
            page_imgs.sort(key=lambda x: x[0], reverse=True)
            for area, raw, ext in page_imgs[:MAX_IMGS_PER_PAGE]:
                data_url, w2, h2, area2 = _b64_data_url(raw, prefer_format=ext)
                out.append({
                    "kind": "image",
                    "data_url": data_url,
                    "page": 1,
                    "width": w2,
                    "height": h2,
                    "area": area2,
                    "caption": None
                })
    except Exception:
        pass

    return out


def _ingest_pptx(path: str, seen_hashes: set[str]) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    prs = Presentation(path)
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text or ""
                if t.strip():
                    texts.append(t)
        if texts:
            out.append({"kind": "text", "text": "\n".join(texts), "page": idx})

        slide_imgs: List[Tuple[int, bytes, str]] = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                try:
                    img = shape.image
                    blob = img.blob
                    hsh = _hash_bytes(blob)
                    if DEDUPLICATE_IMAGES and hsh in seen_hashes:
                        continue
                    with Image.open(io.BytesIO(blob)) as im:
                        w, h = im.size
                        area = w * h
                    if area < MIN_IMAGE_AREA:
                        continue
                    slide_imgs.append((area, blob, (img.ext or "jpeg").lower()))
                    if DEDUPLICATE_IMAGES:
                        seen_hashes.add(hsh)
                except Exception:
                    continue

        if slide_imgs:
            slide_imgs.sort(key=lambda x: x[0], reverse=True)
            for area, raw, ext in slide_imgs[:MAX_IMGS_PER_PAGE]:
                data_url, w2, h2, area2 = _b64_data_url(raw, prefer_format=ext)
                out.append({
                    "kind": "image",
                    "data_url": data_url,
                    "page": idx,
                    "width": w2,
                    "height": h2,
                    "area": area2,
                    "caption": None
                })
    return out
