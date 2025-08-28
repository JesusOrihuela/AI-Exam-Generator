import os
import base64
import fitz
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import io

class DocumentProcessor:
    """
    Maneja la carga de archivos, la extracción de contenido y la conversión
    de imágenes a formatos compatibles con la API.
    """
    SUPPORTED_FORMATS = {'png', 'jpeg', 'gif', 'webp'}

    def _get_supported_image_base64(self, image_bytes, image_format):
        image_format = image_format.lower()
        if image_format == 'jpg': image_format = 'jpeg'

        if image_format in self.SUPPORTED_FORMATS:
            return f"data:image/{image_format};base64,{base64.b64encode(image_bytes).decode('utf-8')}"

        try:
            image = Image.open(io.BytesIO(image_bytes))
            output_buffer = io.BytesIO()
            image.save(output_buffer, format="PNG")
            converted_image_bytes = output_buffer.getvalue()
            return f"data:image/png;base64,{base64.b64encode(converted_image_bytes).decode('utf-8')}"
        except Exception as e:
            print(f"Advertencia: No se pudo convertir una imagen de formato '{image_format}': {e}. Se omitirá.")
            return None

    def process_files(self, file_paths):
        content = []
        for path in file_paths:
            ext = os.path.splitext(path)[1].lower()
            try:
                if ext == ".pdf": content.extend(self._process_pdf(path))
                elif ext == ".docx": content.extend(self._process_docx(path))
                elif ext == ".pptx": content.extend(self._process_pptx(path))
            except Exception as e:
                print(f"Advertencia: No se pudo procesar '{os.path.basename(path)}': {e}. Se omitirá.")
        return content

    def _process_pdf(self, path):
        doc_content = []
        doc = fitz.open(path)
        for page in doc:
            doc_content.append({"type": "text", "text": page.get_text()})
            for img in page.get_images(full=True):
                base_image = doc.extract_image(img[0])
                base64_image = self._get_supported_image_base64(base_image["image"], base_image["ext"])
                if base64_image:
                    doc_content.append({"type": "image_url", "image_url": {"url": base64_image}})
        return doc_content

    def _process_docx(self, path):
        doc_content = []
        doc = DocxDocument(path)
        for para in doc.paragraphs: doc_content.append({"type": "text", "text": para.text})
        for rel in doc.part.rels:
            if "image" in doc.part.rels[rel].target_ref:
                image_part = doc.part.rels[rel].target_part
                ext = image_part.content_type.split('/')[-1]
                base64_image = self._get_supported_image_base64(image_part.blob, ext)
                if base64_image:
                    doc_content.append({"type": "image_url", "image_url": {"url": base64_image}})
        return doc_content

    def _process_pptx(self, path):
        doc_content = []
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): doc_content.append({"type": "text", "text": shape.text})
                if shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
                    image = shape.image
                    base64_image = self._get_supported_image_base64(image.blob, image.ext)
                    if base64_image:
                        doc_content.append({"type": "image_url", "image_url": {"url": base64_image}})
        return doc_content

    def chunk_content(self, content, chunk_size_chars=8000):
        """Divide el contenido extraído en fragmentos de un tamaño aproximado."""
        chunks = []
        current_chunk = []
        current_chunk_len = 0

        for item in content:
            item_len = len(item.get('text', '')) if item['type'] == 'text' else 0
            
            if current_chunk_len + item_len > chunk_size_chars and current_chunk:
                chunks.append(current_chunk)
                current_chunk = []
                current_chunk_len = 0

            current_chunk.append(item)
            current_chunk_len += item_len
        
        if current_chunk:
            chunks.append(current_chunk)
            
        return chunks